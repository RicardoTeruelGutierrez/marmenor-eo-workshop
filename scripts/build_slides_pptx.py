"""Build the three PowerPoint decks (Module 1, 2, 3) for the Mar Menor workshop.

Brand: matches slides/intro.html — ink navy, warm paper, terracotta accent, teal.
Dark cover + closing, light content (sandwich). Speaker notes on every slide.
Run:  python scripts/build_slides_pptx.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = Path(__file__).resolve().parent.parent / "slides"
OUT.mkdir(exist_ok=True)

# ── Palette ────────────────────────────────────────────────────────────────
INK    = RGBColor(0x0E, 0x1A, 0x26)
PAPER  = RGBColor(0xF6, 0xF1, 0xE7)
PAPER2 = RGBColor(0xEE, 0xE5, 0xD4)
ACCENT = RGBColor(0xC2, 0x41, 0x0C)   # terracotta
TEAL   = RGBColor(0x11, 0x5E, 0x59)
GREY   = RGBColor(0x6B, 0x61, 0x57)
INKTX  = RGBColor(0x2A, 0x25, 0x20)
CREAMTX= RGBColor(0xD4, 0xCD, 0xB9)
CARD   = RGBColor(0xFB, 0xF8, 0xF1)

SERIF = "Georgia"
SANS  = "Calibri"
MONO  = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def deck():
    p = Presentation()
    p.slide_width = EMU_W
    p.slide_height = EMU_H
    return p


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def setrun(p, text, size, color, font=SANS, bold=False, italic=False, spacing=None):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.name = font
    r.font.bold = bold
    r.font.italic = italic
    return r


def rect(slide, x, y, w, h, color, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def label(slide, text, x=0.92, y=0.6, color=ACCENT):
    tf = textbox(slide, x, y, 11.5, 0.4)
    setrun(tf.paragraphs[0], text.upper(), 11, color, MONO, bold=True)
    return tf


def footer(slide, left, right):
    tf = textbox(slide, 0.92, 6.95, 11.5, 0.35)
    p = tf.paragraphs[0]
    setrun(p, left.upper(), 8.5, GREY, MONO)
    tf2 = textbox(slide, 0.92, 6.95, 11.5, 0.35)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    setrun(p2, right.upper(), 8.5, GREY, MONO)


def title(slide, lines, x=0.92, y=1.0, size=30, color=INKTX, italic=False, w=11.5):
    tf = textbox(slide, x, y, w, 1.6)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.0
        setrun(p, ln, size, color, SERIF, bold=False, italic=italic)
    return tf


# ── Slide templates ─────────────────────────────────────────────────────────
def cover(prs, mono, big_lines, lede, stats, foot):
    s = blank(prs); bg(s, INK)
    rect(s, 0.92, 1.0, 1.6, 0.09, ACCENT)
    tf = textbox(s, 0.92, 0.62, 11.5, 0.4)
    setrun(tf.paragraphs[0], mono.upper(), 11, ACCENT, MONO, bold=True)
    tt = textbox(s, 0.92, 1.35, 11.5, 2.6)
    for i, ln in enumerate(big_lines):
        p = tt.paragraphs[0] if i == 0 else tt.add_paragraph()
        p.line_spacing = 1.02
        setrun(p, ln, 44, PAPER, SERIF, italic=True)
    ld = textbox(s, 0.92, 4.15, 9.6, 1.4)
    setrun(ld.paragraphs[0], lede, 17, CREAMTX, SANS)
    # stat row
    n = len(stats); gap = 0.45; sw = (11.5 - gap*(n-1)) / n
    for i, (num, key) in enumerate(stats):
        x = 0.92 + i*(sw+gap)
        tnum = textbox(s, x, 5.7, sw, 0.7)
        setrun(tnum.paragraphs[0], num, 30, ACCENT, SERIF, italic=True)
        tkey = textbox(s, x, 6.35, sw, 0.5)
        setrun(tkey.paragraphs[0], key.upper(), 8.5, CREAMTX, MONO)
    footer(s, foot[0], foot[1])
    return s


def closing(prs, mono, big_lines, takeaways, foot):
    s = blank(prs); bg(s, INK)
    rect(s, 0.92, 1.0, 1.6, 0.09, ACCENT)
    tf = textbox(s, 0.92, 0.62, 11.5, 0.4)
    setrun(tf.paragraphs[0], mono.upper(), 11, ACCENT, MONO, bold=True)
    tt = textbox(s, 0.92, 1.3, 11.5, 1.0)
    for i, ln in enumerate(big_lines):
        p = tt.paragraphs[0] if i == 0 else tt.add_paragraph()
        setrun(p, ln, 32, PAPER, SERIF, italic=True)
    y = 2.7
    for i, (head, body) in enumerate(takeaways):
        tn = textbox(s, 0.92, y, 0.6, 0.6)
        setrun(tn.paragraphs[0], str(i+1), 22, ACCENT, MONO, bold=True)
        tb = textbox(s, 1.6, y, 10.8, 0.95)
        p = tb.paragraphs[0]; p.line_spacing = 1.05
        setrun(p, head + "  ", 15, PAPER, SANS, bold=True)
        setrun(p, body, 15, CREAMTX, SANS)
        y += 1.0
    footer(s, foot[0], foot[1])
    return s


def content(prs, lab, ttl_lines, foot, title_size=30, title_italic=False):
    s = blank(prs); bg(s, PAPER)
    label(s, lab)
    title(s, ttl_lines, size=title_size, italic=title_italic)
    footer(s, foot[0], foot[1])
    return s


def card(slide, x, y, w, h, head, body):
    rect(slide, x, y, 0.06, h, TEAL)            # teal left edge
    rect(slide, x+0.06, y, w-0.06, h, CARD)
    th = textbox(slide, x+0.28, y+0.18, w-0.5, 0.4)
    setrun(th.paragraphs[0], head.upper(), 11, TEAL, MONO, bold=True)
    tbf = textbox(slide, x+0.28, y+0.62, w-0.5, h-0.75)
    p = tbf.paragraphs[0]; p.line_spacing = 1.05
    setrun(p, body, 12.5, INKTX, SANS)


def stat_block(slide, x, y, w, num, key):
    tn = textbox(slide, x, y, w, 0.7)
    setrun(tn.paragraphs[0], num, 32, ACCENT, SERIF, italic=True)
    tk = textbox(slide, x, y+0.7, w, 0.5)
    setrun(tk.paragraphs[0], key.upper(), 9, GREY, MONO)


def bullets(slide, x, y, w, h, head, items, body_size=13.5):
    if head:
        th = textbox(slide, x, y, w, 0.4)
        setrun(th.paragraphs[0], head, 15, INKTX, SANS, bold=True)
        y += 0.5
    tf = textbox(slide, x, y, w, h)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.08
        p.space_after = Pt(6)
        setrun(p, "— ", body_size, ACCENT, SANS, bold=True)
        if isinstance(it, tuple):
            setrun(p, it[0] + " ", body_size, INKTX, SANS, bold=True)
            setrun(p, it[1], body_size, INKTX, SANS)
        else:
            setrun(p, it, body_size, INKTX, SANS)


def paragraph(slide, x, y, w, h, text, size=15, color=INKTX, italic=False):
    tf = textbox(slide, x, y, w, h)
    p = tf.paragraphs[0]; p.line_spacing = 1.12
    setrun(p, text, size, color, SANS, italic=italic)


def codebox(slide, x, y, w, h, lines):
    rect(slide, x, y, w, h, INK)
    tf = textbox(slide, x+0.25, y+0.2, w-0.45, h-0.35)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.15
        color = GREY if ln.strip().startswith("#") else CREAMTX
        setrun(p, ln if ln else " ", 11.5, color, MONO)


def simple_table(slide, x, y, w, rows, colw, header=True, fs=12):
    nrows, ncols = len(rows), len(rows[0])
    h = Inches(0.5 * nrows)
    tbl = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y),
                                 Inches(w), h).table
    for ci, cw in enumerate(colw):
        tbl.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER2 if (header and ri == 0) else CARD
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = cell.text_frame.paragraphs[0]
            ish = header and ri == 0
            setrun(para, str(val), 10 if ish else fs,
                   TEAL if ish else INKTX, MONO if ish else SANS, bold=ish)
    return tbl


# ===========================================================================
# MODULE 1
# ===========================================================================
def build_module1():
    prs = deck()
    cover(prs, "Module 01 · Notebook 01_satellite_data_extraction",
          ["Reading chlorophyll", "and heat from orbit."],
          "Sentinel-2 water quality and Sentinel-3 sea-surface temperature over "
          "the Mar Menor — from raw reflectance to the signature of a dying lagoon.",
          [("10 m", "S2 pixel"), ("NDCI", "chl index"),
           ("+2 °C", "2021 heatwave"), ("A·B·C", "three sections")],
          ("module 1 / 00", "satellite data extraction · 75 min"))
    notes(prs.slides[-1],
        "Welcome to Module 1. Goal: go from raw satellite data to the visible "
        "fingerprints of the Mar Menor's collapse. Two sensors: Sentinel-2 for "
        "water colour at 10 m, Sentinel-3 for sea-surface temperature. Stress the "
        "three-section structure (A synthetic/offline, B real workflow, C real "
        "imagery) — it's how the notebook avoids leaving students stuck.")

    s = content(prs, "§ 01 — why the lagoon is hard",
                ["Case-2 waters break the textbook", "ocean-colour algorithms."],
                ("module 1 / 01", "case-2 waters"), title_size=28)
    paragraph(s, 0.92, 2.35, 6.2, 2.2,
        "Open-ocean (Case-1) algorithms assume phytoplankton alone drives the "
        "colour. The Mar Menor is shallow, wind-stirred and nutrient-loaded, so "
        "three things confound the signal at once. The fix is the C2RCC "
        "neural-network correction, then a locally-calibrated index.", size=15)
    paragraph(s, 0.92, 4.7, 6.2, 1.6,
        "Caveat we state out loud: Sentinel-2 L2A surface reflectance is a LAND "
        "correction — not, by itself, an aquatic C2RCC product.", size=13,
        color=ACCENT, italic=True)
    bullets(s, 7.5, 2.3, 4.9, 2.2, "The three confounders", [
        ("Suspended sediment", "raises blue/green → chl-a over-estimated"),
        ("CDOM", "from runoff absorbs in the blue"),
        ("Macrophytes", "on the bottom add a NIR signal")])
    bullets(s, 7.5, 4.75, 4.9, 1.4, "Our chl proxy", [
        ("Rrs(B05)/Rrs(B04)", "— red-edge over red rises with chlorophyll")])
    notes(prs.slides[-1],
        "The single most important conceptual slide. Explain WHY standard ocean "
        "colour fails here: sediment, CDOM and macrophytes each move reflectance "
        "independently of chlorophyll. That's the definition of Case-2 water. "
        "Then the honesty point: L2A (Sen2Cor) is a land atmospheric correction; "
        "for rigorous water work you'd run C2RCC. The proxy B05/B04 exploits the "
        "red-edge: chlorophyll absorbs red (B04), reflects red-edge (B05).")

    s = content(prs, "§ 02 — how the notebook is built",
                ["Three clearly separated sections."],
                ("module 1 / 02", "structure A·B·C"))
    card(s, 0.92, 2.4, 3.7, 2.7, "A · Synthetic",
         "Cached datasets, fully reproducible OFFLINE in class. The 12-point "
         "chl-a series and the SST cube. No internet, no account.")
    card(s, 4.82, 2.4, 3.7, 2.7, "B · Real workflow",
         "The patterns to discover and download real Sentinel data from CDSE / "
         "STAC. Shown, not run in class (needs a free account).")
    card(s, 8.72, 2.4, 3.7, 2.7, "C · Real imagery",
         "Instructor figures from real Sentinel-2 GeoTIFFs: True-Colour, NDCI, "
         "zones. The lagoon outline is the actual satellite coastline.")
    paragraph(s, 0.92, 5.4, 11.5, 0.9,
        "The split keeps “what runs anywhere” cleanly apart from “what needs the "
        "cloud” — so students never get stuck mid-lesson.", size=15, italic=True)
    notes(prs.slides[-1],
        "Explain the A/B/C design — it directly answers a reviewer concern that "
        "synthetic and real data were visually mixed. Section A always runs "
        "(offline). Section B is the real-data recipes (account needed). Section "
        "C renders real downloaded imagery and self-skips if the GeoTIFFs aren't "
        "present. This is the spine of the whole notebook.")

    s = content(prs, "§ 03 — section A · the long view",
                ["Four crises, one chlorophyll curve."],
                ("module 1 / 03", "chlorophyll time series"))
    rows = [["Year", "Event"],
            ["2016", "First green-soup crisis — accumulated nitrate"],
            ["2019", "DANA flood pulse — 300+ mm in 48 h"],
            ["2021", "August anoxia & fish-kill — the climax"],
            ["2025", "Bloom resurgence — early-warning test repeats"]]
    simple_table(s, 0.92, 2.35, 7.0, rows, [1.1, 5.9], fs=12.5)
    bullets(s, 8.3, 2.5, 4.1, 3,
        "What the figure shows", [
        "12 sampling points → lagoon-wide weekly mean",
        "Documented crises light up in the right months",
        "Proof the synthetic data is calibrated, not invented"], body_size=13)
    notes(prs.slides[-1],
        "Section A uses the cached synthetic series. The teaching point: the four "
        "documented crises appear at the correct dates and magnitudes because the "
        "synthetic generator is calibrated against published literature. Walk the "
        "timeline 2016 → 2025. This sets up why 2021 is the focus.")

    s = content(prs, "§ 04 — section A · the kill mechanism",
                ["A marine heatwave that", "strangles the oxygen."],
                ("module 1 / 04", "heatwave → anoxia"), title_size=28)
    stat_block(s, 0.92, 2.5, 3.5, "30.5 °C", "Aug 2021 mean SST")
    stat_block(s, 0.92, 3.9, 3.5, "+1.9 °C", "anomaly vs climatology")
    stat_block(s, 0.92, 5.3, 3.5, "−0.3 mg/L", "O₂ capacity lost")
    bullets(s, 5.0, 2.5, 7.4, 3.4, "The chain, quantified", [
        ("Warmer water holds less O₂", "(Benson–Krause solubility)."),
        ("27.5 → 30.5 °C removes ~0.3 mg/L", "of oxygen capacity."),
        ("At night the bloom consumes O₂", "— and stratification seals the bottom."),
        ("Result:", "bottom anoxia — the 2021 fish-kill condition.")], body_size=14)
    paragraph(s, 5.0, 5.7, 7.4, 0.8,
        "Teaching point: SST is not a curiosity — it is a computable driver of "
        "the ecological collapse.", size=13, color=TEAL, italic=True)
    notes(prs.slides[-1],
        "The mechanism slide. Make the causal chain explicit and quantitative: a "
        "+2 °C heatwave lowers O₂ solubility by ~0.3 mg/L; the bloom respires at "
        "night; stratification prevents reoxygenation of the bottom. Together → "
        "hypoxia → fish-kill. We show the SST field clipped to the real lagoon "
        "outline plus the Benson-Krause O₂ solubility curve.")

    s = content(prs, "§ 05 — section B · how to get the data",
                ["The current Copernicus access path."],
                ("module 1 / 05", "CDSE · STAC · COG"))
    codebox(s, 0.92, 2.35, 6.0, 3.4, [
        "# Current CDSE STAC endpoint (Nov 2025)",
        "from pystac_client import Client",
        "",
        'cat = Client.open(',
        '  "https://stac.dataspace.copernicus.eu/v1/")',
        "",
        "search = cat.search(",
        '  collections=["sentinel-2-l2a"],',
        '  datetime="2021-08-01/2021-09-30",',
        "  bbox=BBOX,",
        '  query={"eo:cloud_cover": {"lt": 20}})'])
    bullets(s, 7.3, 2.5, 5.1, 3.4, "What students must remember", [
        "SciHub is dead; CDSE is the entry point",
        ("New STAC root", "stac.dataspace.copernicus.eu/v1"),
        ("Collection is lowercase", "sentinel-2-l2a"),
        ("Use L2A", "(surface), not L1C (top-of-atmosphere)"),
        ("The COG trick:", "read only the AOI window, a few MB")], body_size=12.5)
    notes(prs.slides[-1],
        "Section B is the real-world recipe — shown, not run in class. The "
        "critical correction: the old catalogue.dataspace.copernicus.eu/stac was "
        "deprecated Nov 2025; the current root is stac.dataspace.copernicus.eu/v1 "
        "and the collection id is lowercase sentinel-2-l2a. Emphasise L2A vs L1C "
        "and the cloud-optimised-GeoTIFF window-read trick.")

    s = content(prs, "§ 06 — section C · real imagery",
                ["Pixel-perfect, from real scenes."],
                ("module 1 / 06", "true colour · NDCI · zones"))
    card(s, 0.92, 2.4, 5.65, 1.85, "True Colour",
         "Four cloud-free, full-coverage scenes (Feb 2020 → Mar 2022), chosen by "
         "measuring real coverage AND cloud over the lagoon — not the catalogue %.")
    card(s, 6.77, 2.4, 5.65, 1.85, "NDCI per pixel",
         "Chlorophyll index over 1.3 M water pixels. Clear winter ≈ −0.28; summer "
         "2021 bloom ≈ −0.05 — a large jump on a shared colour scale.")
    card(s, 0.92, 4.45, 5.65, 1.85, "Zone analysis",
         "North / Centre / South medians reveal the gradient driven by "
         "agricultural runoff on the north-west shore.")
    card(s, 6.77, 4.45, 5.65, 1.85, "Outline from NIR",
         "The lagoon boundary is extracted from the near-infrared band itself — "
         "191 vertices, exact, not a hand-drawn polygon.")
    notes(prs.slides[-1],
        "Section C is the visual payoff with real Sentinel-2 pixels. Four key "
        "products. Stress two methodological wins: (1) scene selection measures "
        "coverage AND cloud over the lagoon — that's why the 14-Jul-2021 "
        "swath-edge scene was rejected; (2) the lagoon outline is derived from the "
        "NIR band, so borders are pixel-perfect.")

    closing(prs, "§ 07 — what to walk away with",
        ["Module 1 in four sentences."], [
        ("CDSE, not SciHub.", "STAC at stac.dataspace.copernicus.eu/v1, collection sentinel-2-l2a."),
        ("Case-2 means be careful.", "Atmospheric-correct with C2RCC; never trust raw L2A as an aquatic product."),
        ("NDCI separates clear water from bloom", "and exposes the north–south nutrient gradient."),
        ("A +2 °C heatwave removes ~0.3 mg/L O₂", "— enough, with night-time respiration, to drive the 2021 anoxia.")],
        ("module 1 / 07", "fin · module 1"))
    notes(prs.slides[-1], "Recap and bridge to Module 2: now that we can read the "
        "lagoon from space, can we validate it against the buoys and predict the "
        "crises with ML?")
    prs.save(OUT / "module1.pptx")
    print("module1.pptx:", len(list(prs.slides)), "slides")


# ===========================================================================
# MODULE 2
# ===========================================================================
def build_module2():
    prs = deck()
    cover(prs, "Module 02 · Notebook 02_insitu_timeseries_ml",
          ["From buoys to", "early warning."],
          "Integrating the in-situ network with the satellite record, validating "
          "retrievals, and learning to predict and flag the lagoon's crises.",
          [("5", "buoy stations"), ("18k", "daily records"),
           ("2", "honest CV schemes"), ("42 d", "best lead time")],
          ("module 2 / 00", "in-situ + machine learning · 70 min"))
    notes(prs.slides[-1],
        "Module 2 connects satellite to ground truth and adds machine learning. "
        "Two big honesty themes to flag up front: (1) validate spatially, not just "
        "in time; (2) an early-warning detector must be judged out-of-sample. "
        "These reflect the review fixes that made the module defensible.")

    s = content(prs, "§ 01 — where the ground truth lives",
                ["The regional monitoring networks."],
                ("module 2 / 01", "in-situ networks"))
    card(s, 0.92, 2.4, 5.65, 1.85, "CARM Mar Menor",
         "Real-time sondes around the lagoon: T, S, chl-a, DO, turbidity. The "
         "dashboard regional authorities actually watch.")
    card(s, 6.77, 2.4, 5.65, 1.85, "IMIDA / SIAM",
         "Agro-meteorological stations: soil moisture, nitrate proxies — the "
         "upstream pressure on the lagoon.")
    card(s, 0.92, 4.45, 5.65, 1.85, "CHS-SAIH",
         "Segura watershed hydrology: precipitation and runoff at 15-minute "
         "cadence — the flood signal.")
    card(s, 6.77, 4.45, 5.65, 1.85, "Copernicus Marine INS-TAC",
         "Aggregated in-situ ocean observations (buoys, CTDs, gliders) for the "
         "Mediterranean.")
    notes(prs.slides[-1],
        "Map the data landscape for in-situ. Our cached buoy file mimics CARM "
        "exports: 5 stations, daily, 2016-2025, with T, S, chl-a, turbidity, DO, "
        "nitrate. Point out that the upstream (IMIDA/CHS) is what ultimately "
        "drives the lagoon — the nutrient pressure.")

    s = content(prs, "§ 02 — the crises in the raw signal",
                ["What the agency saw in real time."],
                ("module 2 / 02", "hypoxia · DANA"))
    bullets(s, 0.92, 2.5, 5.6, 2.6, "August 2021 — the hypoxia", [
        ("Dissolved oxygen crashes below 2 mg/L", "at the worst stations"),
        ("SST peaks simultaneously", "— the heatwave"),
        ("This is the plot behind the fish-kill headlines", "")], body_size=14)
    bullets(s, 6.7, 2.5, 5.7, 2.6, "September 2019 — the DANA", [
        ("300+ mm in 48 h", "cut-off cold drop"),
        ("Salinity drops ~4 PSU", "as freshwater floods in"),
        ("Nitrate spikes ten-fold", "— the watershed plume")], body_size=14)
    paragraph(s, 0.92, 5.5, 11.5, 0.8,
        "Before any modelling, we simply PLOT the events — the data must show the "
        "story we claim to study.", size=15, italic=True)
    notes(prs.slides[-1],
        "Exploratory plots first. Show the 2021 DO collapse against the 2 mg/L "
        "hypoxia line with SST overlaid, and the 2019 DANA as a precip → nitrate "
        "→ salinity cascade. Lesson: always eyeball the events before trusting any "
        "model on them.")

    s = content(prs, "§ 03 — validation",
                ["Does the satellite agree", "with the water?"],
                ("module 2 / 03", "satellite ↔ in-situ match-ups"), title_size=28)
    paragraph(s, 0.92, 2.5, 6.0, 2.4,
        "We co-locate each Sentinel-2 sampling point with its nearest buoy "
        "(within 4 km, same day) and score the retrieval the way an ocean-colour "
        "paper would. A scatter against the 1:1 line is the honest visual test.",
        size=15)
    bullets(s, 7.3, 2.5, 5.1, 2.6, "The metrics that matter", [
        ("R²", "— variance explained"),
        ("RMSE / RMSLE", "— error (log form: chl-a is log-normal)"),
        ("Bias", "— systematic over/under-estimate"),
        ("1:1 line", "on the scatter — the visual truth test")], body_size=13)
    notes(prs.slides[-1],
        "Match-ups: the discipline of validation. 3 km / 3 h is the ocean-colour "
        "convention; we use 4 km / same-day because in-situ is daily. Report R², "
        "RMSE, RMSLE (log because chl-a is log-normal) and bias, and always draw "
        "the 1:1 line. Never claim skill without a match-up.")

    s = content(prs, "§ 04 — machine learning, done honestly",
                ["Two cross-validations,", "two different questions."],
                ("module 2 / 04", "gradient boosting · honest CV"), title_size=28)
    rows = [["Scheme", "Question it answers", "Result"],
            ["TimeSeriesSplit", "Predict future dates? (time order)", "high R²"],
            ["GroupKFold by station", "Predict an unseen station? (space)", "R² ≈ 0.17"]]
    simple_table(s, 0.92, 2.7, 11.5, rows, [3.0, 6.0, 2.5], fs=12.5)
    paragraph(s, 0.92, 4.6, 11.5, 1.4,
        "The gap is the lesson: a model that looks great in time can be weak in "
        "space. Report the harder number. A gradient-boosting regressor on "
        "reflectances + band ratios + season is the workhorse.", size=15,
        italic=True)
    notes(prs.slides[-1],
        "The key ML-rigor slide. TimeSeriesSplit answers 'can I predict future "
        "dates' and looks strong. GroupKFold holding out whole stations answers "
        "'can I predict an unmonitored place' and is much harder (R² around 0.17). "
        "The honest deployment number is the spatial one. This directly addresses "
        "the reviewer's 'too optimistic ML' concern.")

    s = content(prs, "§ 05 — multivariate anomaly detection",
                ["An early-warning detector", "that never peeks at the future."],
                ("module 2 / 05", "isolation forest · out-of-sample"), title_size=27)
    codebox(s, 0.92, 2.55, 6.0, 3.3, [
        "# Fit ONLY on quiet years 2017-18",
        "train = sm.query(\"date < '2019-01-01'\")",
        "iso = IsolationForest(",
        "    n_estimators=300,",
        "    contamination=0.03,",
        "    random_state=42)",
        "iso.fit(train[FEATURES])",
        "",
        "# Apply out-of-sample (2019+)",
        "sm['anomaly'] = iso.predict(sm[FEATURES])"])
    bullets(s, 7.3, 2.6, 5.1, 3.2, "Why quiet years only", [
        "Fit on the whole series and it learns the crises as 'normal'",
        "Train on 2017–2018, apply to 2019+ → every flag is out-of-sample",
        "Watches chl-a, turbidity, DO, SST, salinity, nitrate jointly"],
        body_size=13)
    notes(prs.slides[-1],
        "Anomaly detection with honesty built in. If you fit IsolationForest on "
        "all years, the crises become part of 'normal' and the early-warning "
        "claim is circular. We fit only on the quiet 2017-18 baseline and apply "
        "forward. It's multivariate — joint anomalies across six variables catch "
        "precursors a single chl-a threshold misses.")

    s = content(prs, "§ 06 — the metric agencies actually want",
                ["Lead time ≠ detection delay."],
                ("module 2 / 06", "lead time vs delay"))
    rows = [["Crisis", "First flag", "Outcome"],
            ["2019 DANA", "−42 days", "genuine early warning"],
            ["2021 anoxia", "+2 days", "caught after onset (a miss)"],
            ["2025 bloom", "+1 day", "caught after onset (a miss)"]]
    simple_table(s, 0.92, 2.5, 11.5, rows, [2.8, 2.7, 6.0], fs=12.5)
    paragraph(s, 0.92, 4.6, 11.5, 1.5,
        "We separate a flag BEFORE the start (lead time — the win) from a flag "
        "AFTER (detection delay — a miss). Lumping them together is how monitoring "
        "papers quietly oversell themselves.", size=15, italic=True)
    notes(prs.slides[-1],
        "The subtle but important fix. The naive code scanned a window starting "
        "before and ending after the event, took the first flag, and called it "
        "'lead time' — even when it came after onset. We split lead time (before) "
        "from detection delay (after). Honest result: 2019 was a real 42-day "
        "warning; 2021 and 2025 were caught only after the crisis began.")

    closing(prs, "§ 07 — what to walk away with",
        ["Module 2 in four sentences."], [
        ("Validate before you model.", "Match-ups with R²/RMSE/bias and a 1:1 line keep you honest."),
        ("Time CV and space CV differ.", "Report the unseen-station score, not just the easy one."),
        ("Early-warning trains on quiet years", "and must be judged out-of-sample."),
        ("Lead time and detection delay are opposites.", "Only a flag before onset is a real warning.")],
        ("module 2 / 07", "fin · module 2"))
    notes(prs.slides[-1], "Recap. Bridge to Module 3: Modules 1-2 used data we "
        "prepared — Module 3 teaches students to build that data layer themselves.")
    prs.save(OUT / "module2.pptx")
    print("module2.pptx saved")


# ===========================================================================
# MODULE 3
# ===========================================================================
def build_module3():
    prs = deck()
    cover(prs, "Module 03 · Notebook 03_build_database",
          ["Build your own", "data pipeline."],
          "From zero to a working database: discover, download live, clean, store "
          "in SQLite + Parquet, and turn it into a pollution study.",
          [("8", "free sources"), ("3 %", "of a scene pulled"),
           ("SQL", "+ Parquet store"), ("0", "images for a series")],
          ("module 3 / 00", "from zero to pro · 90 min"))
    notes(prs.slides[-1],
        "Module 3 is the capstone: students build the whole data layer themselves "
        "and run real downloads live. Frame it with the kitchen metaphor (grocery "
        "run → wash → chop → fridge → cook → restaurant). Everything runs in class "
        "with no account; the Copernicus account section then shows what a free "
        "login unlocks.")

    s = content(prs, "§ 01 — the shape of the whole module",
                ["Six steps, one pipeline."],
                ("module 3 / 01", "pipeline overview"))
    steps = [("1 · grocery", "Discover & download", "free archives"),
             ("2 · wash", "Clean & QA", "drop clouds, land"),
             ("3 · chop", "Extract", "tidy rows"),
             ("4 · fridge", "Store", "SQLite + Parquet"),
             ("5 · cook", "Analyse", "pollution study")]
    n = len(steps); gap = 0.3; w = (11.5 - gap*(n-1)) / n
    for i, (k, h, sub) in enumerate(steps):
        x = 0.92 + i*(w+gap)
        rect(s, x, 2.6, w, 0.06, TEAL)
        rect(s, x, 2.66, w, 1.9, CARD)
        tk = textbox(s, x+0.15, 2.8, w-0.3, 0.4)
        setrun(tk.paragraphs[0], k.upper(), 8.5, ACCENT, MONO, bold=True)
        th = textbox(s, x+0.15, 3.2, w-0.3, 0.8)
        ph = th.paragraphs[0]; ph.line_spacing = 0.95
        setrun(ph, h, 14, INKTX, SERIF, bold=True)
        tsub = textbox(s, x+0.15, 4.0, w-0.3, 0.5)
        setrun(tsub.paragraphs[0], sub, 11, GREY, SANS)
        if i < n-1:
            ar = textbox(s, x+w-0.02, 3.2, 0.34, 0.5)
            setrun(ar.paragraphs[0], "→", 18, ACCENT, MONO)
    paragraph(s, 0.92, 5.2, 11.5, 1.0,
        "The kitchen metaphor sticks: a project is a grocery run, washing, "
        "chopping, a fridge, and cooking — then Step 6 is running the restaurant "
        "every day (an incremental, scheduled job).", size=15, italic=True)
    notes(prs.slides[-1],
        "Give students the whole map before details. Walk the six steps with the "
        "kitchen analogy — it makes an abstract pipeline memorable. Each step maps "
        "to a section of the notebook. Step 6 (productionise) is what turns a "
        "one-off analysis into infrastructure.")

    s = content(prs, "§ 02 — where the free data is",
                ["Know the archive, know the login."],
                ("module 3 / 02", "data landscape"))
    rows = [["Source", "What you get", "Login?", "Best for"],
            ["NASA GIBS", "Daily true-colour, 2000→now", "none", "browse any date"],
            ["Earth Search → AWS COG", "Sentinel-2 bands, 10 m", "none", "analysis pixels"],
            ["Open-Meteo Air-Quality", "PM2.5, PM10, ozone", "none", "pollution context"],
            ["Open-Meteo Marine", "Sea-surface temperature", "none", "in-water conditions"],
            ["Copernicus CDSE / SH", "On-demand C2RCC, any product", "free acct", "real-time processing"],
            ["CMEMS · Earthdata · EMODnet", "Reanalysis, ocean colour, bathy", "free/OGC", "archives & context"]]
    simple_table(s, 0.92, 2.35, 11.5, rows, [3.0, 4.0, 1.5, 3.0], fs=11)
    notes(prs.slides[-1],
        "The data landscape. Rule of thumb: prototype today with the four no-login "
        "sources (GIBS, AWS COG, Open-Meteo air-quality and marine); graduate to "
        "account-based ones (CDSE, CMEMS, Earthdata) when the project needs a "
        "specific product or the full archive. Note Open-Meteo air quality is "
        "directly relevant to pollution.")

    s = content(prs, "§ 03 — downloads you run in the room",
                ["Real data, on screen, no login."],
                ("module 3 / 03", "live, no credentials"))
    card(s, 0.92, 2.4, 5.65, 1.85, "GIBS image",
         "One URL → a true-colour PNG of the lagoon for any date. Change the "
         "date, re-run, browse.")
    card(s, 6.77, 2.4, 5.65, 1.85, "Sentinel-2 COG window",
         "Search → open the COG header → read only the AOI window: ~8 MB instead "
         "of ~240 MB (3 % of the scene).")
    card(s, 0.92, 4.45, 5.65, 1.85, "Air-quality series",
         "GET → JSON → DataFrame: hourly PM2.5 / PM10 / ozone at the lagoon.")
    card(s, 6.77, 4.45, 5.65, 1.85, "Marine SST series",
         "Same pattern, live sea-surface temperature — a free complement to "
         "Sentinel-3.")
    notes(prs.slides[-1],
        "The 'see the bytes arrive' moment — the heart of the user's request. All "
        "four cells run live with no credentials and degrade gracefully offline. "
        "The COG window read is the headline skill: pull 3 % of a scene over HTTP. "
        "Encourage students to change the date and their own coordinates and "
        "re-run.")

    s = content(prs, "§ 04 — what a free Copernicus account unlocks",
                ["Server-side processing,", "not just raw bands."],
                ("module 3 / 04", "CDSE process + statistical API"), title_size=28)
    codebox(s, 0.92, 2.55, 6.0, 3.4, [
        "# Send a recipe, get the product",
        'EVALSCRIPT = """//VERSION=3',
        'function setup(){return {',
        '  input:["B04","B05","dataMask"],',
        '  output:{bands:1}};}',
        'function evaluatePixel(s){',
        '  return [(s.B05-s.B04)/',
        '          (s.B05+s.B04)];}"""',
        "",
        "# Process API  -> NDCI map, any date",
        "# Statistical API -> series, 0 downloads"])
    bullets(s, 7.3, 2.6, 5.1, 3.4, "The three unlocks", [
        ("Process API", "— send an evalscript, get the finished product"),
        ("Statistical API", "— a full AOI time series, zero imagery downloaded"),
        ("Raw bands", "— read the chlorophyll signature off the spectrum")],
        body_size=12.5)
    notes(prs.slides[-1],
        "The account section — verified live against CDSE. Registration takes ~2 "
        "min; store the OAuth client id/secret as env vars and the same cells run "
        "live. The big unlock vs the AWS mirror: server-side processing. Process "
        "API returns a finished NDCI/true-colour/band product for any date; "
        "Statistical API returns a whole time series with NO image downloads — the "
        "cheap way to run a monitoring dashboard.")

    s = content(prs, "§ 05 — store it the right way",
                ["SQLite for queries,", "Parquet for bulk."],
                ("module 3 / 05", "schema · idempotent"), title_size=28)
    paragraph(s, 0.92, 2.55, 5.7, 2.6,
        "CSVs become a swamp. A database gives types, primary keys (no "
        "duplicates), indices and SQL. We use a pragmatic hybrid: SQLite for the "
        "catalogue + tidy observations + stations; Parquet for the millions of "
        "per-pixel values. Idempotent writes (INSERT OR REPLACE on a UNIQUE key) "
        "mean re-running never duplicates — essential for a scheduled job.",
        size=14)
    codebox(s, 6.9, 2.55, 5.5, 3.0, [
        "CREATE TABLE observations (",
        "  obs_id   INTEGER PRIMARY KEY,",
        "  date     TEXT, source TEXT,",
        "  variable TEXT, zone   TEXT,",
        "  value    REAL, units  TEXT,",
        "  UNIQUE(date,source,variable,zone)",
        ");",
        "-- tidy, long format = easy SQL"])
    notes(prs.slides[-1],
        "Why a database, not CSVs: types, keys, indices, SQL. The hybrid: SQLite "
        "for queryable tables, Parquet for bulk arrays. The single most important "
        "engineering idea here is idempotency — the UNIQUE constraint plus INSERT "
        "OR REPLACE means a scheduled pipeline can re-process any date without "
        "creating duplicates. Tidy long format makes every later query a one-liner.")

    s = content(prs, "§ 06 — the payoff",
                ["A pollution study, straight from SQL."],
                ("module 3 / 06", "trend · exceedance · join"))
    bullets(s, 0.92, 2.5, 6.4, 3.2, "One query each", [
        "Monthly-mean chl-a series",
        ("Robust long-term trend", "(Theil-Sen)"),
        ("'Poor status' exceedance days", "per year"),
        "North / Centre / South zone comparison",
        ("Satellite ↔ in-situ JOIN", "match-up")], body_size=13.5)
    stat_block(s, 7.8, 2.7, 4.5, "20 %", "2021 obs above threshold")
    stat_block(s, 7.8, 4.1, 4.5, "8.2 %", "2019 obs above threshold")
    paragraph(s, 7.8, 5.3, 4.6, 1.0,
        "The database detects the documented crises by itself — proof the pipeline "
        "works end to end.", size=13, color=TEAL, italic=True)
    notes(prs.slides[-1],
        "The payoff: with everything in SQL, monitoring questions become short "
        "queries — monthly means, a robust Theil-Sen trend, threshold-exceedance "
        "counts, zone comparisons, and a satellite-vs-in-situ JOIN. The validation "
        "that it works: the exceedance counts spike in exactly 2019 and 2021, the "
        "documented crisis years.")

    closing(prs, "§ 07 — scale & takeaways",
        ["Make it grow on its own."], [
        ("Discover → subset → array/table", "is the same pattern for every source."),
        ("Move as few pixels as possible.", "COG windows + server-side Statistical API."),
        ("Clean before you store; store tidy.", "Idempotent SQLite + Parquet scales to a scheduled job."),
        ("The SQL transfers", "to PostGIS / DuckDB. Swap the AOI and the database builds itself.")],
        ("module 3 / 07", "fin · module 3"))
    notes(prs.slides[-1], "Close on empowerment: they now own the entire chain and "
        "ran it live. Scaling path: PostGIS for spatial/multi-user, DuckDB for "
        "querying Parquet at scale, cron/Airflow for scheduling. The whole thing "
        "re-points to any study area by changing the AOI.")
    prs.save(OUT / "module3.pptx")
    print("module3.pptx saved")


if __name__ == "__main__":
    build_module1()
    build_module2()
    build_module3()
    print("All three decks written to", OUT)
