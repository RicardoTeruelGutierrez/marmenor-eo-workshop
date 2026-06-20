# -*- coding: utf-8 -*-
"""Build the Live Session PowerPoint — slides/session.pptx.

An editable deck mirroring slides/session.html: same story, agenda, figure
slides (with the real session figures) and takeaways, in the workshop identity.
Run:  python scripts/build_session_pptx.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
FIGS = ROOT / "slides" / "figs"
OUT  = ROOT / "slides" / "session.pptx"

INK   = RGBColor(0x0D, 0x1B, 0x2A)
INK2  = RGBColor(0x16, 0x26, 0x3A)
PAPER = RGBColor(0xF6, 0xF1, 0xE7)
ACCENT= RGBColor(0xC2, 0x41, 0x0C)
TEAL  = RGBColor(0x2A, 0x9D, 0x8F)
GOLD  = RGBColor(0xE9, 0xA4, 0x4B)
GREY  = RGBColor(0x6B, 0x61, 0x57)
CREAM = RGBColor(0xCB, 0xD5, 0xDD)
INKTX = RGBColor(0x2A, 0x25, 0x20)
SERIF = "Georgia"; MONO = "Consolas"; SANS = "Calibri"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = INK if dark else PAPER
    return s

def box(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tf = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"):
        setattr(tf, m, 0)
    return tf

def run(p, text, size, color, font=SANS, bold=False, italic=False):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = font
    r.font.bold = bold; r.font.italic = italic
    return r

def kick(s, text, color=ACCENT):
    tf = box(s, 0.9, 0.55, 11.5, 0.4)
    run(tf.paragraphs[0], text.upper(), 11, color, MONO, bold=True)

def foot(s, num, dark=False):
    tf = box(s, 0.9, 7.0, 11.5, 0.35)
    run(tf.paragraphs[0], str(num), 8.5, (CREAM if dark else GREY), MONO)
    tf2 = box(s, 0.9, 7.0, 11.5, 0.35); p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    run(p, "MAR MENOR · LIVE EARTH-OBSERVATION SESSION", 8.5, (CREAM if dark else GREY), MONO)

def title(s, lines, color=INKTX, size=30, y=1.0, italic=False, w=11.5):
    tf = box(s, 0.9, y, w, 1.7)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.02
        run(p, ln, size, color, SERIF, italic=italic)

def fig(s, name, x, y, w, h):
    p = FIGS / f"sess_{name}.png"
    if not p.exists(): return
    # dark card behind the image
    card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = INK; card.line.fill.background()
    card.shadow.inherit = False
    from PIL import Image
    iw, ih = Image.open(p).size
    pad = 0.12; aw, ah = w - 2*pad, h - 2*pad
    sc = min(aw/iw, ah/ih); dw, dh = iw*sc, ih*sc
    s.shapes.add_picture(str(p), Inches(x+(w-dw)/2), Inches(y+(h-dh)/2),
                         Inches(dw), Inches(dh))


# ── Cover ────────────────────────────────────────────────────────────────
s = slide(dark=True)
kick(s, "Live Earth-Observation Session · 60 minutes", GOLD)
title(s, ["Watch a lagoon", "collapse, from space."], PAPER, 40, 1.3, italic=True)
tf = box(s, 0.9, 3.5, 9.4, 1.5)
run(tf.paragraphs[0], "In one hour you download real satellite and pollution data "
    "with your own hands — and see the 2021 Mar Menor crisis appear on screen. "
    "No install, mostly no account.", 16, CREAM, SANS)
stats = [("10 m","Sentinel-2 detail"),("5","live data sources"),
         ("7","hands-on steps"),("0","software to install")]
for i,(n,k) in enumerate(stats):
    x = 0.9 + i*3.0
    t1 = box(s, x, 5.7, 2.8, 0.7); run(t1.paragraphs[0], n, 30, GOLD, SERIF, italic=True)
    t2 = box(s, x, 6.35, 2.8, 0.5); run(t2.paragraphs[0], k.upper(), 8.5, CREAM, MONO)
foot(s, "00", dark=True)

# ── Story ────────────────────────────────────────────────────────────────
s = slide(dark=True)
kick(s, "The place · the crisis", GOLD)
title(s, ["A small sea that broke, on camera."], PAPER, 26, 1.0, w=7.0)
tf = box(s, 0.9, 2.2, 5.6, 4.2)
p = tf.paragraphs[0]; run(p, "The Mar Menor is Europe's largest coastal lagoon — "
    "shallow, closed, ringed by farmland.", 15, CREAM, SANS)
p2 = tf.add_paragraph(); p2.space_before = Pt(10)
run(p2, "In August 2021 it turned green and three tonnes of fish died: an algal "
    "bloom, fed by farm runoff and a marine heatwave, stripped the water of oxygen.",
    15, CREAM, SANS)
p3 = tf.add_paragraph(); p3.space_before = Pt(10)
run(p3, "Every stage of it was recorded from orbit — free, open, and the subject "
    "of today's session.", 13, GREY, SANS, italic=True)
fig(s, "truecolor", 6.9, 1.9, 5.5, 4.6)
foot(s, "01", dark=True)

# ── Agenda ───────────────────────────────────────────────────────────────
s = slide()
kick(s, "What we'll do together")
title(s, ["Seven steps, one hour."], INKTX, 30, 1.0)
rows = [("1","Browse any date from orbit — NASA GIBS"),
        ("2","Download real Sentinel-2 (10 m) — only our window"),
        ("3","Measure the bloom — the NDCI chlorophyll index"),
        ("4","Add live pollution data — air quality"),
        ("5","Unlock server products with a free Copernicus key"),
        ("6","Detect the crises with machine learning"),
        ("7","Store & query it all in a database")]
y = 2.35
for i, t in rows:
    ti = box(s, 0.9, y, 0.6, 0.5); run(ti.paragraphs[0], i, 18, ACCENT, MONO, bold=True)
    tt = box(s, 1.7, y, 10.6, 0.5); run(tt.paragraphs[0], t, 15.5, INKTX, SANS)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(y+0.5),
                            Inches(11.5), Pt(0.75))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xD8,0xCF,0xBC); ln.line.fill.background()
    ln.shadow.inherit = False
    y += 0.62
foot(s, "02")


def figure_slide(num, kicktext, ttl, name, tag, blurb, bullets=None, wide=False, formula=None):
    s = slide()
    kick(s, kicktext)
    title(s, [ttl], INKTX, 28, 0.95)
    if wide:
        tf = box(s, 0.9, 2.0, 11.5, 1.5)
        if tag: run(tf.paragraphs[0], tag.upper()+"   ", 10, TEAL, MONO, bold=True)
        run(tf.paragraphs[0], blurb, 15, INKTX, SANS)
        if bullets:
            for b in bullets:
                pb = tf.add_paragraph(); pb.space_before = Pt(4)
                run(pb, "— "+b, 13, GREY, SANS)
        fig(s, name, 0.9, 3.5, 11.5, 3.0)
    else:
        tf = box(s, 0.9, 2.2, 5.4, 4.2, MSO_ANCHOR.TOP)
        if tag: run(tf.paragraphs[0], tag.upper()+"   ", 10, TEAL, MONO, bold=True)
        run(tf.paragraphs[0], blurb, 16, INKTX, SANS)
        if formula:
            pf = tf.add_paragraph(); pf.space_before = Pt(10)
            run(pf, formula, 14, ACCENT, MONO)
        if bullets:
            for b in bullets:
                pb = tf.add_paragraph(); pb.space_before = Pt(8)
                run(pb, "— "+b, 13.5, GREY, SANS)
        fig(s, name, 6.5, 1.9, 5.9, 4.6)
    foot(s, num)

figure_slide("03", "Step 1 · browse", "Find a clear day from orbit.", "gibs",
             "live · no login",
             "NASA GIBS returns a true-colour picture of any day since 2000. We scout "
             "for cloud-free dates before downloading the heavy data.",
             bullets=["One URL, one image — change the date and re-run."])
figure_slide("04", "Step 2 · download", "What each wavelength sees.", "bands",
             "live · no login",
             "Sentinel-2 records separate bands. We read only our window (~16 MB, not "
             "800). Watch the bloom flip from dark in red to bright in red-edge.", wide=True)
figure_slide("05", "Step 3 · measure", "The bloom, as a number.", "ndci", "live",
             "The NDCI turns two bands into chlorophyll at every pixel. Land is greyed "
             "out; the bloom is strongest in the north-west, by the farmland.",
             formula="NDCI = (B05 − B04) / (B05 + B04)")
figure_slide("06", "Step 4 · pollution", "More than water colour.", "airquality",
             "live · no login",
             "A keyless API returns hourly PM2.5, PM10 and ozone — the dust over the same "
             "farmland that feeds the lagoon. Same pattern: ask, receive, plot.")
figure_slide("07", "Step 5 · go pro", "Five products, one request.", "wq_panel",
             "free account",
             "With a free Copernicus key the server computes products for you. One call "
             "returns the three optical components of the water — plus the bloom peak.",
             bullets=["True colour · chlorophyll · bloom hotspots · turbidity · CDOM."], wide=True)
figure_slide("08", "Step 6 · machine learning", "It finds the crises itself.", "anomaly",
             "live",
             "An Isolation Forest, trained only on the calm years, flags every later day "
             "that looks abnormal — and rediscovers the 2019 and 2021 crises it never saw.",
             bullets=["100% of the documented crisis days, caught out-of-sample."])

# ── Takeaways ────────────────────────────────────────────────────────────
s = slide(dark=True)
kick(s, "What you take away", GOLD)
title(s, ["The barrier was never access."], PAPER, 28, 1.0)
items = [("Free & open","the data that recorded a disaster is one request away — no paywall."),
         ("Move few pixels","search a catalogue, read only your window, let the server compute."),
         ("One pattern","discover → take a subset → analyse → store. It works for any coast."),
         ("You did it live","browse, download, measure, detect and store — in one hour, no install.")]
y = 2.5
for k,(h,t) in enumerate(items):
    ti = box(s, 0.9, y, 0.6, 0.5); run(ti.paragraphs[0], str(k+1), 20, GOLD, MONO, bold=True)
    tt = box(s, 1.6, y, 10.8, 0.8); p = tt.paragraphs[0]
    run(p, h+".  ", 15, PAPER, SANS, bold=True); run(p, t, 15, CREAM, SANS)
    y += 0.95
tf = box(s, 0.9, 6.4, 11.5, 0.4)
run(tf.paragraphs[0], "Go further → notebooks 01 · 02 · 03 in the repo", 10, GREY, MONO)
foot(s, "fin", dark=True)

prs.save(OUT)
print(f"session.pptx written — {len(prs.slides._sldIdLst)} slides")
